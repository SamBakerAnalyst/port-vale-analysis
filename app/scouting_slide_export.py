from __future__ import annotations

from datetime import datetime
from io import BytesIO
from typing import Any, Iterable, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.pdf_report import pdf_safe
from app.slide_export import _add_rect, _add_textbox, _new_widescreen_presentation

WHITE = (255, 255, 255)
INK = (17, 24, 39)
MUTED = (107, 114, 128)
LINE = (229, 231, 235)
STRIPE = (249, 250, 251)
HEAD = (243, 244, 246)
NAVY = (30, 58, 95)

FONT = "Calibri"
ROWS_PER_SLIDE = 20
HEADERS = ("#", "Player", "Age", "Min", "Club", "League", "Ovr")


def _rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def _chunks(items: Sequence[Any], size: int) -> Iterable[Sequence[Any]]:
    for index in range(0, len(items), size):
        yield items[index : index + size]


def _cell(
    cell: Any,
    text: str,
    *,
    bold: bool = False,
    size: int = 12,
    fill: tuple[int, int, int] = WHITE,
    color: tuple[int, int, int] = INK,
    align: int = PP_ALIGN.LEFT,
) -> None:
    cell.text = pdf_safe(text)
    frame = cell.text_frame
    frame.word_wrap = False
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = Pt(6)
    frame.margin_right = Pt(6)
    p = frame.paragraphs[0]
    p.alignment = align
    p.font.name = FONT
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color)
    cell.fill.solid()
    cell.fill.fore_color.rgb = _rgb(fill)


def _cover(prs: Presentation, body: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w = prs.slide_width
    h = prs.slide_height
    _add_rect(slide, 0, 0, w, h, WHITE)
    _add_rect(slide, 0, 0, w, Inches(0.06), NAVY)

    _add_textbox(slide, Inches(0.6), Inches(0.45), w - Inches(1.2), Inches(0.3), "Impect scouting", size=11, color=MUTED)
    _add_textbox(
        slide,
        Inches(0.6),
        Inches(0.85),
        w - Inches(1.2),
        Inches(0.7),
        pdf_safe(body.position_label or "Long list"),
        size=32,
        bold=True,
        color=INK,
    )

    leagues = ", ".join(body.leagues or [])
    bits = [
        f"Leagues: {leagues}" if leagues else "",
        f"{int(body.min_minutes):d}+ minutes",
        f"{len(body.players or [])} players",
        "Scores vs peers in same league",
    ]
    y = Inches(1.75)
    for bit in [b for b in bits if b]:
        _add_textbox(slide, Inches(0.6), y, w - Inches(1.2), Inches(0.28), bit, size=13, color=MUTED)
        y += Inches(0.32)

    when = body.generated_at or datetime.now().strftime("%d %b %Y")
    _add_textbox(slide, Inches(0.6), h - Inches(0.55), w - Inches(1.2), Inches(0.25), f"Generated {when}", size=10, color=MUTED)


def _list_slide(
    prs: Presentation,
    *,
    position_label: str,
    page_players: Sequence[Any],
    page_num: int,
    total_pages: int,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    w = prs.slide_width
    _add_rect(slide, 0, 0, w, prs.slide_height, WHITE)

    title = f"{position_label}  ·  {page_players[0].rank}–{page_players[-1].rank}  ·  {page_num}/{total_pages}"
    _add_textbox(slide, Inches(0.45), Inches(0.22), w - Inches(0.9), Inches(0.28), pdf_safe(title), size=13, bold=True, color=INK)

    rows = [[
        str(p.rank),
        p.name or "",
        "" if p.age is None else str(p.age),
        "" if p.minutes is None else str(int(round(p.minutes))),
        p.club or "—",
        p.league or "—",
        "—" if p.overall is None else f"{p.overall:.1f}",
    ] for p in page_players]

    n_rows = len(rows) + 1
    n_cols = len(HEADERS)
    left = Inches(0.45)
    top = Inches(0.58)
    width = w - Inches(0.9)
    row_h = Inches(0.3)
    height = row_h * n_rows

    table = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table

    col_widths = [Inches(0.38), Inches(2.35), Inches(0.42), Inches(0.52), Inches(2.05), Inches(1.15), Inches(0.48)]
    for i, cw in enumerate(col_widths):
        table.columns[i].width = int(cw)

    for col, header in enumerate(HEADERS):
        _cell(table.cell(0, col), header, bold=True, size=11, fill=HEAD, color=INK,
              align=PP_ALIGN.CENTER if col != 1 else PP_ALIGN.LEFT)

    for r, row in enumerate(rows, start=1):
        fill = WHITE if r % 2 else STRIPE
        for c, value in enumerate(row):
            is_name = c == 1
            is_ovr = c == 6
            _cell(
                table.cell(r, c),
                value,
                bold=is_ovr,
                size=12,
                fill=fill,
                color=NAVY if is_ovr else INK,
                align=PP_ALIGN.LEFT if is_name else PP_ALIGN.CENTER,
            )


def build_scouting_slides_pptx(body: Any) -> bytes:
    players = body.players or []
    if not players:
        raise ValueError("No players to export.")

    pages = list(_chunks(players, ROWS_PER_SLIDE))
    prs = _new_widescreen_presentation()
    _cover(prs, body)

    label = body.position_label or "Long list"
    for i, page in enumerate(pages, start=1):
        _list_slide(prs, position_label=label, page_players=page, page_num=i, total_pages=len(pages))

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()

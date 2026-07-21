from __future__ import annotations

from datetime import datetime
from io import BytesIO
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from app.label_utils import humanize_profile_name
from app.metric_bars import format_metric_value, format_percentile_label, percentile_colors
from app.pdf_report import decode_image_data, pdf_safe

PHOTO_PANEL_WIDTH = Inches(4.35)


def _rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(*color)


def _add_rect(slide: Any, left: float, top: float, width: float, height: float, fill: tuple[int, int, int]) -> None:
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.fill.background()


def _add_textbox(
    slide: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: tuple[int, int, int] = (15, 23, 42),
    align: int = PP_ALIGN.LEFT,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = pdf_safe(text)
    paragraph.alignment = align
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _rgb(color)


def _add_photo_panel(slide: Any, prs: Presentation, top: float, height: float, dark: bool = False) -> None:
    left = prs.slide_width - PHOTO_PANEL_WIDTH - Inches(0.45)
    fill = (30, 41, 59) if dark else (241, 245, 249)
    _add_rect(slide, left, top, PHOTO_PANEL_WIDTH, height, fill)
    _add_textbox(
        slide,
        left,
        top + (height / 2) - Inches(0.15),
        PHOTO_PANEL_WIDTH,
        Inches(0.35),
        "Player image",
        size=12,
        color=(203, 213, 225) if dark else (100, 116, 139),
        align=PP_ALIGN.CENTER,
    )


def _left_content_width(prs: Presentation) -> float:
    return prs.slide_width - PHOTO_PANEL_WIDTH - Inches(1.15)


def _add_factor_bar_row(
    slide: Any,
    left: float,
    top: float,
    width: float,
    label: str,
    percentile: float | None,
    raw_value: float | None,
) -> None:
    label_w = Inches(1.65)
    badge_w = Inches(0.42)
    gap = Inches(0.08)
    track_w = width - label_w - badge_w - (gap * 2)
    bar_h = Inches(0.22)

    _add_textbox(slide, left, top, label_w, bar_h, label, size=11, color=(100, 116, 139))

    track_x = left + label_w + gap
    _add_rect(slide, track_x, top, track_w, bar_h, (232, 237, 244))

    fill_color, badge_bg, badge_text = percentile_colors(percentile)
    if percentile is not None:
        fill_w = max(track_w * (percentile / 100.0), Inches(0.35))
        _add_rect(slide, track_x, top, fill_w, bar_h, fill_color)
        if raw_value is not None:
            _add_textbox(
                slide,
                track_x,
                top + Inches(0.03),
                fill_w,
                bar_h,
                format_metric_value(raw_value),
                size=9,
                bold=True,
                color=(255, 255, 255),
                align=PP_ALIGN.CENTER,
            )

    badge_x = track_x + track_w + gap
    _add_rect(slide, badge_x, top, badge_w, bar_h, badge_bg)
    _add_textbox(
        slide,
        badge_x,
        top + Inches(0.03),
        badge_w,
        bar_h,
        format_percentile_label(percentile),
        size=9,
        bold=True,
        color=badge_text,
        align=PP_ALIGN.CENTER,
    )


def _add_metric_grid(slide: Any, entry: Any, left: float, top: float, width: float) -> None:
    players = entry.players or []
    labels = entry.labels or []
    if not players or not labels:
        return

    row_h = Inches(0.28)
    cursor_y = top
    multi_player = len(players) > 1

    for index, label in enumerate(labels):
        for player_index, player in enumerate(players):
            percentile = player.radar_values[index] if index < len(player.radar_values) else None
            raw_values = getattr(player, "raw_values", None) or []
            raw_value = raw_values[index] if index < len(raw_values) else None
            row_label = label if not multi_player or player_index == 0 else f"  {player.player.split(' ')[0]}"
            _add_factor_bar_row(slide, left, cursor_y, width, row_label, percentile, raw_value)
            cursor_y += row_h


def _cover_slide(prs: Presentation, body: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    width = prs.slide_width
    height = prs.slide_height
    left_width = _left_content_width(prs)

    _add_rect(slide, 0, 0, width, height, (15, 23, 42))
    _add_rect(slide, 0, height - Inches(0.08), width, Inches(0.08), (56, 189, 248))
    _add_photo_panel(slide, prs, Inches(0.45), height - Inches(0.9), dark=True)

    primary = body.players[0] if body.players else None
    player_name = pdf_safe(primary.player if primary else "Player report")

    _add_textbox(slide, Inches(0.55), Inches(0.55), left_width, Inches(0.4), "PORT VALE FC", size=14, bold=True, color=(56, 189, 248))
    _add_textbox(slide, Inches(0.55), Inches(1.0), left_width, Inches(1.2), player_name, size=38, bold=True, color=(255, 255, 255))

    y = Inches(2.15)
    if primary:
        meta_parts = [part for part in [primary.season_label, primary.position_label] if part]
        if meta_parts:
            _add_textbox(slide, Inches(0.55), y, left_width, Inches(0.45), "  |  ".join(meta_parts), size=15, color=(203, 213, 225))
            y += Inches(0.45)

    if body.benchmark_subtitle:
        _add_textbox(slide, Inches(0.55), y, left_width, Inches(0.8), body.benchmark_subtitle, size=12, color=(148, 163, 184))
        y += Inches(0.5)
    if body.profiles:
        profile_labels = [humanize_profile_name(name) for name in body.profiles]
        _add_textbox(
            slide,
            Inches(0.55),
            y,
            left_width,
            Inches(1.0),
            "Profiles: " + ", ".join(profile_labels),
            size=12,
            color=(148, 163, 184),
        )

    _add_textbox(
        slide,
        Inches(0.55),
        height - Inches(0.55),
        left_width,
        Inches(0.3),
        pdf_safe(f"Generated {body.generated_at or datetime.now().strftime('%d %b %Y, %H:%M')}"),
        size=10,
        color=(148, 163, 184),
    )


def _chart_slide(prs: Presentation, title: str, subtitle: str, image_data: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    width = prs.slide_width
    left_width = _left_content_width(prs)
    _add_rect(slide, 0, 0, width, Inches(0.9), (15, 23, 42))
    _add_rect(slide, 0, Inches(0.9), width, Inches(0.05), (56, 189, 248))
    _add_textbox(slide, Inches(0.55), Inches(0.18), left_width, Inches(0.45), title, size=24, bold=True, color=(255, 255, 255))
    if subtitle:
        _add_textbox(slide, Inches(0.55), Inches(0.55), left_width, Inches(0.3), subtitle, size=12, color=(203, 213, 225))
    _add_photo_panel(slide, prs, Inches(1.0), prs.slide_height - Inches(1.15))

    image_bytes = decode_image_data(image_data)
    slide.shapes.add_picture(BytesIO(image_bytes), Inches(0.55), Inches(1.05), width=left_width, height=Inches(5.9))


def _profile_slide(prs: Presentation, title: str, subtitle: str, image_data: str, entry: Any) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    width = prs.slide_width
    left_width = _left_content_width(prs)
    _add_rect(slide, 0, 0, width, Inches(0.9), (15, 23, 42))
    _add_rect(slide, 0, Inches(0.9), width, Inches(0.05), (56, 189, 248))
    _add_textbox(slide, Inches(0.55), Inches(0.18), left_width, Inches(0.45), title, size=22, bold=True, color=(255, 255, 255))
    if subtitle:
        _add_textbox(slide, Inches(0.55), Inches(0.55), left_width, Inches(0.3), subtitle, size=12, color=(203, 213, 225))
    _add_photo_panel(slide, prs, Inches(1.0), prs.slide_height - Inches(1.15))

    factor_count = len(entry.labels or [])
    if factor_count > 6:
        chart_height = Inches(3.0)
    elif factor_count > 4:
        chart_height = Inches(3.35)
    else:
        chart_height = Inches(4.1)
    image_bytes = decode_image_data(image_data)
    slide.shapes.add_picture(BytesIO(image_bytes), Inches(0.55), Inches(1.0), width=left_width, height=chart_height)
    _add_metric_grid(slide, entry, Inches(0.55), Inches(1.0) + chart_height + Inches(0.12), left_width)


def _full_bleed_slide(prs: Presentation, image_data: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    image_bytes = decode_image_data(image_data)
    slide.shapes.add_picture(
        BytesIO(image_bytes),
        0,
        0,
        width=prs.slide_width,
        height=prs.slide_height,
    )


def _new_widescreen_presentation() -> Presentation:
    prs = Presentation()
    # Standard 16:9 widescreen (13.333" × 7.5") — not the default 4:3 template size.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def build_whole_deck_pptx(body: Any) -> bytes:
    if not body.sections:
        raise ValueError("No chart images to export.")

    prs = _new_widescreen_presentation()
    for section in body.sections:
        _full_bleed_slide(prs, section.image_data)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def build_coach_slides_pptx(body: Any) -> bytes:
    if not body.sections:
        raise ValueError("No chart images to export.")

    export_mode = getattr(body, "export_mode", "coach") or "coach"
    if export_mode == "whole":
        return build_whole_deck_pptx(body)

    prs = _new_widescreen_presentation()

    _cover_slide(prs, body)
    section_map = {section.title: section for section in body.sections}

    radar = section_map.get("Profile radar")
    if radar is not None:
        _chart_slide(prs, "Profile overview", "Percentile scores vs the cross-league benchmark cohort.", radar.image_data)

    for entry in body.drilldowns:
        section = section_map.get(entry.profile)
        if section is None:
            continue
        subtitle = f"{len(entry.labels)} factors assessed"
        if len(entry.players) > 1:
            subtitle += f"  |  {len(entry.players)} players compared"
        _profile_slide(prs, humanize_profile_name(entry.profile), subtitle, section.image_data, entry)

    pizza = section_map.get("Squad percentile pizza")
    if pizza is not None:
        _chart_slide(prs, "Squad percentile view", "Distribution view against the benchmark cohort.", pizza.image_data)

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()

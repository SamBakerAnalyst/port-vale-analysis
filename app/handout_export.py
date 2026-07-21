from __future__ import annotations

from io import BytesIO
from typing import Any

from fpdf import FPDF
from pptx import Presentation
from pptx.util import Mm

from app.pdf_report import decode_image_data

# True A4 portrait (matches West Brom handout reference)
A4_WIDTH_MM = 210.0
A4_HEIGHT_MM = 297.0


class HandoutA4PDF(FPDF):
    """A4 portrait — one full-bleed handout slide per page."""

    def __init__(self) -> None:
        super().__init__(orientation="P", unit="mm", format="A4")
        self.set_auto_page_break(auto=False)
        self.set_margins(0, 0, 0)

    def add_page_image(self, image_data: str) -> None:
        self.add_page()
        self.set_fill_color(255, 255, 255)
        self.rect(0, 0, A4_WIDTH_MM, A4_HEIGHT_MM, style="F")
        self.image(
            BytesIO(decode_image_data(image_data)),
            x=0,
            y=0,
            w=A4_WIDTH_MM,
            h=A4_HEIGHT_MM,
        )


def _iter_page_images(pages: list[Any]) -> list[str]:
    images: list[str] = []
    for page in pages:
        if isinstance(page, dict):
            image_data = page.get("imageData") or page.get("image_data") or ""
        else:
            image_data = getattr(page, "image_data", None) or getattr(page, "imageData", "")
        if image_data:
            images.append(image_data)
    return images


def build_handout_export_pdf(body: Any) -> bytes:
    pages = getattr(body, "pages", None) or []
    images = _iter_page_images(list(pages))
    if not images:
        raise ValueError("No export pages to render.")

    pdf = HandoutA4PDF()
    for image_data in images:
        pdf.add_page_image(image_data)

    buffer = BytesIO()
    pdf.output(buffer)
    return buffer.getvalue()


def _new_a4_portrait_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Mm(A4_WIDTH_MM)
    prs.slide_height = Mm(A4_HEIGHT_MM)
    return prs


def build_handout_export_pptx(body: Any) -> bytes:
    pages = getattr(body, "pages", None) or []
    images = _iter_page_images(list(pages))
    if not images:
        raise ValueError("No export pages to render.")

    prs = _new_a4_portrait_presentation()
    blank = prs.slide_layouts[6]
    for image_data in images:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            BytesIO(decode_image_data(image_data)),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
